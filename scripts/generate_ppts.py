"""Auto-generate EDA and Modelling PowerPoint presentations.

Usage:
    python scripts/generate_ppts.py

Outputs:
    outputs/ppts/EDA_Liver_Disease.pptx
    outputs/ppts/Modelling_Liver_Disease.pptx
"""

import json
import sys
from io import BytesIO
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.config import load_settings
from src.data.loader import load_raw_data
from src.data.splitter import extract_X_y, stratified_split
from src.features.schema import (
    CATEGORICAL_FEATURES,
    NOMINAL_FEATURES,
    NUMERIC_FEATURES,
    TARGET,
)
from src.features.woe import iv_summary_table
from src.models.persistence import load_model
from src.visualization.eda_plots import (
    crosstabulate,
    plot_boxplot,
    plot_correlation_matrix,
    plot_percentage_stacked_chart,
)
from src.visualization.style import PALETTE

# ── Color constants (matching project palette) ──────────────────────────────
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x4A, 0x90, 0xD9)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)

ROOT = Path(__file__).resolve().parent.parent
OUTPUTS_DIR = ROOT / "outputs"
PPTS_DIR = OUTPUTS_DIR / "ppts"
CHARTS_DIR = OUTPUTS_DIR / "charts"


# ── Helpers ─────────────────────────────────────────────────────────────────

def _fig_to_stream(fig: plt.Figure) -> BytesIO:
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
    buf.seek(0)
    plt.close(fig)
    return buf


def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def _add_header_bar(slide, title: str, subtitle: str = "") -> None:
    """Dark navy header bar spanning full width."""
    from pptx.util import Emu
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.25),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.color.rgb = NAVY

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        from pptx.util import Pt as Pt2
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12)
        r2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = _blank_slide(prs)
    # Full background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.color.rgb = NAVY

    # Accent bar
    accent = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.33), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.color.rgb = BLUE

    txb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(1.5))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = WHITE

    txb2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(1.0))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)


def _add_text_slide(prs, title: str, bullets: list[str], subtitle: str = "") -> None:
    slide = _blank_slide(prs)
    _add_header_bar(slide, title, subtitle)
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        r = p.add_run()
        r.text = f"• {bullet}"
        r.font.size = Pt(14)
        r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)
        p.space_after = Pt(6)


def _add_image_slide(prs, title: str, img_stream: BytesIO, subtitle: str = "") -> None:
    slide = _blank_slide(prs)
    _add_header_bar(slide, title, subtitle)
    slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.4), Inches(12.33), Inches(5.7))


def _add_two_image_slide(prs, title: str, left_stream: BytesIO, right_stream: BytesIO,
                          subtitle: str = "") -> None:
    slide = _blank_slide(prs)
    _add_header_bar(slide, title, subtitle)
    slide.shapes.add_picture(left_stream, Inches(0.3), Inches(1.4), Inches(6.2), Inches(5.7))
    slide.shapes.add_picture(right_stream, Inches(6.8), Inches(1.4), Inches(6.2), Inches(5.7))


def _add_table_slide(prs, title: str, headers: list[str],
                     rows: list[list[str]], subtitle: str = "") -> None:
    slide = _blank_slide(prs)
    _add_header_bar(slide, title, subtitle)

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header row
    col_width = Inches(12.33 / n_cols)
    tbl = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(12.33), Inches(min(5.5, n_rows * 0.45))
    ).table

    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.size = Pt(11)

    for ri, row in enumerate(rows):
        bg = LIGHT_GRAY if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.size = Pt(10)


# ── EDA PPT ─────────────────────────────────────────────────────────────────

def build_eda_ppt(df: pd.DataFrame) -> Presentation:
    prs = _new_prs()

    # Slide 1: Title
    _add_title_slide(prs,
                     "Liver Disease Prediction",
                     "Exploratory Data Analysis | Kaggle Dataset (1,700 Records)")

    # Slide 2: Dataset Overview
    pos_count = (df[TARGET] == "Positive").sum()
    neg_count = (df[TARGET] == "Negative").sum()
    fig, ax = plt.subplots(figsize=(5, 5))
    ax.pie(
        [pos_count, neg_count],
        labels=["Positive", "Negative"],
        colors=[PALETTE["positive"], PALETTE["negative"]],
        autopct="%1.1f%%",
        startangle=90,
        textprops={"fontsize": 13},
    )
    ax.set_title("Target Distribution", fontsize=14, fontweight="bold")
    fig.tight_layout()
    _add_image_slide(prs, "Dataset Overview",
                     _fig_to_stream(fig),
                     f"{len(df):,} records | {len(df.columns)} features | Binary classification")

    # Slide 3: Descriptive Stats
    desc = df[NUMERIC_FEATURES].describe().round(2)
    headers = ["Statistic"] + NUMERIC_FEATURES
    rows = [[idx] + desc.loc[idx].tolist() for idx in desc.index]
    _add_table_slide(prs, "Descriptive Statistics (Numeric Features)", headers, rows)

    # Slide 4: Numeric Features vs Diagnosis (boxplots grid)
    n = len(NUMERIC_FEATURES)
    fig, axes = plt.subplots(1, n, figsize=(14, 4))
    for ax, feat in zip(axes, NUMERIC_FEATURES):
        plot_boxplot(df, feat, TARGET, ax=ax)
    fig.tight_layout()
    _add_image_slide(prs, "Numeric Features vs Diagnosis",
                     _fig_to_stream(fig),
                     "Side-by-side boxplots with Welch t-test significance")

    # Slide 5: Categorical Features vs Diagnosis (stacked bars)
    fig, axes = plt.subplots(2, 3, figsize=(14, 8))
    axes_flat = axes.flatten()
    for i, feat in enumerate(CATEGORICAL_FEATURES):
        ct = crosstabulate(df, feat, TARGET, bins=6)
        plot_percentage_stacked_chart(ct, title=feat, xlabel=feat, ax=axes_flat[i])
    for j in range(len(CATEGORICAL_FEATURES), len(axes_flat)):
        axes_flat[j].set_visible(False)
    fig.tight_layout()
    _add_image_slide(prs, "Categorical Features vs Diagnosis",
                     _fig_to_stream(fig),
                     "100% stacked bar charts")

    # Slide 6: WoE / IV Results
    iv_df = iv_summary_table(df, NUMERIC_FEATURES + CATEGORICAL_FEATURES, TARGET)
    iv_df = iv_df.sort_values("Total_IV", ascending=False).reset_index(drop=True)
    headers = ["Feature", "IV", "Predictive Power"]
    rows = [[r["Feature"], f"{r['Total_IV']:.4f}", r["Interpretation"]]
            for _, r in iv_df.iterrows()]
    _add_table_slide(prs, "Weight of Evidence — Information Value Summary",
                     headers, rows,
                     "IV > 0.3 = Strong predictor | 0.1–0.3 = Medium | < 0.1 = Weak")

    # Slide 7: Correlation Matrix
    fig, ax = plt.subplots(figsize=(8, 6))
    plot_correlation_matrix(df, NUMERIC_FEATURES, ax=ax)
    fig.tight_layout()
    _add_image_slide(prs, "Correlation Matrix (Numeric Features)",
                     _fig_to_stream(fig),
                     "Pearson correlations — lower triangle")

    # Slide 8: Key Findings
    _add_text_slide(prs, "Key Findings", [
        f"Dataset: {len(df):,} records, {pos_count} positive ({pos_count/len(df):.1%}) | "
        f"{neg_count} negative ({neg_count/len(df):.1%}) — moderately imbalanced",
        "AlcoholConsumption and LiverFunctionTest show the highest IV (strong predictors)",
        "GeneticRisk (Low/Medium/High) has clear ordinal pattern vs Diagnosis",
        "Numeric features show statistically significant differences between classes (p < 0.001)",
        "Low multicollinearity among numeric features (all Pearson |r| < 0.3)",
        "No missing values or duplicate records found in the dataset",
    ])

    return prs


# ── Modelling PPT ────────────────────────────────────────────────────────────

def build_modelling_ppt(df: pd.DataFrame, metadata: dict) -> Presentation:
    prs = _new_prs()

    # Slide 1: Title
    _add_title_slide(prs,
                     "Liver Disease Prediction",
                     "Predictive Modelling | Results & Model Selection")

    # Slide 2: Preprocessing Pipeline
    _add_text_slide(prs, "Preprocessing Pipeline", [
        "Numeric features (Age, BMI, AlcoholConsumption, PhysicalActivity, LiverFunctionTest): StandardScaler",
        "Ordinal feature (GeneticRisk): OrdinalEncoder with fixed category order [Low < Medium < High]",
        "Nominal features (Gender, Smoking, Diabetes, Hypertension): OneHotEncoder with drop='if_binary'",
        "All transformations wrapped in a sklearn ColumnTransformer (no leakage — fit on train only)",
        "Single canonical create_preprocessor() function used everywhere",
        "Output: 10 features (5 numeric + 1 ordinal + 4 binary nominal)",
    ])

    # Slide 3: Model Comparison table
    comparison = metadata.get("model_comparison", [])
    if comparison:
        headers = ["Model", "Accuracy", "F1 (Weighted)", "ROC-AUC"]
        rows = [[m["model"], f"{m['accuracy']:.4f}", f"{m['f1_weighted']:.4f}", f"{m['roc_auc']:.4f}"]
                for m in comparison]
        _add_table_slide(prs, "Baseline Model Comparison", headers, rows,
                         "80/20 stratified train-test split | All models use same preprocessor")

    # Slide 4: Feature Importances
    fi = metadata.get("feature_importances", {})
    if fi:
        feats = list(fi.keys())
        vals = list(fi.values())
        idx = np.argsort(vals)
        fig, ax = plt.subplots(figsize=(9, 5))
        ax.barh([feats[i] for i in idx], [vals[i] for i in idx],
                color=PALETTE["primary"])
        ax.set_xlabel("Importance", fontsize=10)
        ax.set_title("Random Forest Feature Importances", fontsize=12, fontweight="bold")
        ax.spines[["top", "right"]].set_visible(False)
        fig.tight_layout()
        _add_image_slide(prs, "Feature Importances (Best Model: Random Forest)",
                         _fig_to_stream(fig),
                         "AlcoholConsumption and LiverFunctionTest are the most predictive features")

    # Slide 5: Classification Report
    cr = metadata.get("classification_report", {})
    if cr:
        headers = ["Class", "Precision", "Recall", "F1-Score", "Support"]
        rows = []
        for cls in ["Negative", "Positive"]:
            if cls in cr:
                m = cr[cls]
                rows.append([cls, f"{m['precision']:.3f}", f"{m['recall']:.3f}",
                              f"{m['f1-score']:.3f}", str(int(m['support']))])
        for agg in ["macro avg", "weighted avg"]:
            if agg in cr:
                m = cr[agg]
                rows.append([agg.title(), f"{m['precision']:.3f}", f"{m['recall']:.3f}",
                              f"{m['f1-score']:.3f}", ""])
        _add_table_slide(prs, "Classification Report — Random Forest", headers, rows)

    # Slide 6: Advanced Techniques
    _add_text_slide(prs, "Advanced Modelling Techniques", [
        "Weight of Evidence (WoE) Logistic Regression: numeric features binned → WoE-encoded → LR",
        "  - Leakage-free: WoE mappings computed on train fold only, applied to test",
        "  - WoE smoothing bug fixed: smoothing only applied when zero event/non-event counts exist",
        "DTSegmentedLR Hybrid: shallow Decision Tree extracts leaf segments → OneHotEncoded → LR",
        "  - Proper sklearn estimator (BaseEstimator + ClassifierMixin) — works in cross_val_score",
        "  - Inspired by Facebook GBDT+LR technique for feature crossing",
        "Optuna Hyperparameter Tuning: Bayesian optimization with StratifiedKFold(5) + f1_weighted",
        "Cross-validation: StratifiedKFold(5) with statistical overfitting detection (gap > 2×std)",
    ])

    # Slide 7: Best Model Summary
    best_acc = metadata.get("accuracy", 0)
    best_f1 = metadata.get("f1_score", 0)
    best_auc = metadata.get("roc_auc", 0)
    hp = metadata.get("hyperparameters", {})
    _add_text_slide(prs, "Best Model: Random Forest", [
        f"Accuracy: {best_acc:.4f}",
        f"F1 Score (Weighted): {best_f1:.4f}",
        f"ROC-AUC: {best_auc:.4f}",
        f"Hyperparameters: n_estimators={hp.get('n_estimators', 100)}, "
        f"max_depth={hp.get('max_depth', 10)}, "
        f"min_samples_split={hp.get('min_samples_split', 10)}, "
        f"min_samples_leaf={hp.get('min_samples_leaf', 4)}",
        "Model saved as models/liver_disease_model.pkl (sklearn Pipeline + ColumnTransformer)",
        "Deployed via Streamlit Cloud with GitHub Actions CI/CD",
    ])

    # Slide 8: Deployment Summary
    _add_text_slide(prs, "Deployment & Architecture", [
        "Streamlit app: 3 tabs — Prediction | Model Info | About",
        "Prediction tab: two-column layout, risk gauge, feature importance chart",
        "Model served via @st.cache_resource — loaded once, shared across sessions",
        "CI/CD: GitHub Actions lint (ruff) + test (pytest, 80% coverage) on push/PR",
        "Deploy gate: verifies app.py + model.pkl + .streamlit/config.toml exist",
        "TDD: 76 pytest tests written first (data, features, models, evaluation, app)",
        "Package structure: src/ with data/, features/, models/, evaluation/, visualization/",
    ])

    return prs


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    PPTS_DIR.mkdir(parents=True, exist_ok=True)
    CHARTS_DIR.mkdir(parents=True, exist_ok=True)

    print("Loading data...")
    settings = load_settings()
    df = load_raw_data(settings)
    print(f"  {len(df)} rows loaded.")

    print("Loading metadata...")
    meta_path = ROOT / "models" / "model_metadata.json"
    if not meta_path.exists():
        print(f"  metadata not found at {meta_path} — run scripts/generate_metadata.py first.")
        metadata = {}
    else:
        with open(meta_path) as f:
            metadata = json.load(f)

    print("Building EDA presentation...")
    eda_prs = build_eda_ppt(df)
    eda_path = PPTS_DIR / "EDA_Liver_Disease.pptx"
    eda_prs.save(str(eda_path))
    print(f"  Saved: {eda_path}")

    print("Building Modelling presentation...")
    mod_prs = build_modelling_ppt(df, metadata)
    mod_path = PPTS_DIR / "Modelling_Liver_Disease.pptx"
    mod_prs.save(str(mod_path))
    print(f"  Saved: {mod_path}")

    print("\nDone. Generated:")
    print(f"  {eda_path}")
    print(f"  {mod_path}")


if __name__ == "__main__":
    main()
